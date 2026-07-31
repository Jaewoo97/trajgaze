"""Redraw a qualitative figure as an editable PowerPoint slide.

Same input as `render_qual_vector.py` (a `--dump-layout` JSON) and the same layout formulas,
so the slide is the figure, not an approximation. The difference is what the output is made
of: every card, token box, gaze ring, chip and string is a native PowerPoint shape or text
box, so the question can be reworded, a row moved, a colour changed or a frame swapped
straight in PowerPoint or Keynote. Only the video frames are images.

  /opt/conda/bin/python scripts/viz_qual/render_qual_pptx.py \
      scripts/viz_qual/layout/sg_O1V0P1F0_idx83_*.json \
      --font scripts/viz_qual/Inter.ttf --out fig_qual_sg83.pptx [--type-scale 1.6]

The slide is Wc x Hc points, one design unit per point, matching the PDF page exactly.

Notes on editing:
- Token boxes of one frame are grouped with their frame, gaze ring and time chip (24 groups
  for a 4-row, 6-frame figure); double-click to get inside, or pass --no-group for 600+ loose
  shapes.
- Text is set in Inter. If Inter is not installed on the machine opening the file, PowerPoint
  substitutes a default sans and the line breaks move; --font-name picks another family.
- Font metrics for wrapping and block heights still come from the .ttf given by --font.
"""
from __future__ import annotations

import argparse
import json
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from layout_common import (C_GAZE, eyebrow, footer,      # noqa: E402
                           legend_entries, resolve_strip)

# design tokens, shared with the PNG and PDF renderers
INK, INK_80, INK_48 = (29, 29, 31), (51, 51, 51), (122, 122, 122)
HAIRLINE, WHITE = (224, 224, 224), (255, 255, 255)
BLUE, OK_FG, BAD_FG = (0, 102, 204), (26, 127, 55), (193, 58, 52)

BOLD_WEIGHTS = {"SemiBold", "Bold", "ExtraBold", "Black"}


class Type:
    """PIL metrics for layout; PowerPoint gets the same sizes in points."""

    def __init__(self, path, scale=1.0):
        self.path, self.scale = path, scale
        self._pil = {}

    def pil(self, size, weight="Regular"):
        key = (size, weight)
        if key not in self._pil:
            f = ImageFont.truetype(self.path, int(round(size * self.scale)))
            try:
                f.set_variation_by_name(weight)
            except Exception:
                pass
            self._pil[key] = f
        return self._pil[key]


def wrap(d, text, font, maxw):
    out = []
    for para in text.split("\n"):
        cur = ""
        for w in para.split(" "):
            t = (cur + " " + w).strip()
            if d.textlength(t, font=font) <= maxw or not cur:
                cur = t
            else:
                out.append(cur)
                cur = w
        out.append(cur)
    return out or [""]


def _srgb(el, tag):
    """The a:srgbClr under spPr/<tag>, or None."""
    sub = el.find(qn(tag))
    return None if sub is None else sub.find(qn("a:srgbClr"))


def _set_alpha(clr, alpha):
    """PowerPoint carries transparency on the colour, and python-pptx has no API for it."""
    if clr is None or alpha >= 1.0:
        return
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(round(alpha * 100000)))})
    clr.append(a)


class Slide:
    """Thin drawing layer over python-pptx, in design units (1 unit = 1 pt)."""

    def __init__(self, prs, tf: Type, font_name: str):
        self.prs, self.T, self.font_name = prs, tf, font_name
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = self.slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*WHITE)
        self.shapes = self.slide.shapes

    # ── shapes ──
    def _host(self, host):
        """`host or self.shapes` would be wrong: a shape tree is a sequence, so an empty group
        is falsy and its first shape would land on the slide instead, above the frame."""
        return self.shapes if host is None else host

    def _shape(self, host, kind, x0, y0, x1, y1, fill=None, line=None, lw=1.0,
               fill_alpha=1.0, line_alpha=1.0, radius=None, name=None):
        w, h = max(x1 - x0, 0.5), max(y1 - y0, 0.5)
        shp = host.add_shape(kind, Pt(x0), Pt(y0), Pt(w), Pt(h))
        shp.shadow.inherit = False
        if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
            shp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(*fill)
            _set_alpha(_srgb(shp._element.spPr, "a:solidFill"), fill_alpha)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = RGBColor(*line)
            shp.line.width = Pt(lw)
            ln = shp._element.spPr.find(qn("a:ln"))
            _set_alpha(None if ln is None else _srgb(ln, "a:solidFill"), line_alpha)
        tf = shp.text_frame
        tf.word_wrap = False
        for side in ("left", "right", "top", "bottom"):
            setattr(tf, f"margin_{side}", 0)
        if name:
            shp.name = name
        return shp

    def rrect(self, x0, y0, x1, y1, r, **kw):
        return self._shape(self._host(kw.pop("host", None)), MSO_SHAPE.ROUNDED_RECTANGLE,
                           x0, y0, x1, y1, radius=r, **kw)

    def rect(self, x0, y0, x1, y1, **kw):
        return self._shape(self._host(kw.pop("host", None)), MSO_SHAPE.RECTANGLE,
                           x0, y0, x1, y1, **kw)

    def ring(self, cx, cy, r, color, lw, alpha=1.0, host=None, name=None):
        return self._shape(self._host(host), MSO_SHAPE.OVAL,
                           cx - r, cy - r, cx + r, cy + r,
                           fill=None, line=color, lw=lw, line_alpha=alpha, name=name)

    def image(self, path, x, y, w, h, host=None, name=None):
        pic = self._host(host).add_picture(path, Pt(x), Pt(y), Pt(w), Pt(h))
        if name:
            pic.name = name
        return pic

    def group(self, name=None):
        grp = self.shapes.add_group_shape()
        if name:
            grp.name = name
        return grp

    # ── text ──
    def text(self, x, y, lines, size, weight="Regular", color=INK, width=None,
             line_h=None, va="top", host=None, name=None, wrap_text=True):
        """`lines` is the already-wrapped list; `y` is the top of the first line's ascent
        (va="top") or the vertical centre of the block (va="center"), as in the PDF path.

        The box is centred on the block and anchored middle, which lands the same way in
        PowerPoint, Keynote and LibreOffice; top-anchoring would depend on each renderer's
        first-baseline rule.
        """
        if isinstance(lines, str):
            lines = [lines]
        f = self.T.pil(size, weight)
        asc, desc = f.getmetrics()
        lh = line_h if line_h is not None else asc + desc
        n = len(lines)
        cy = y + (asc + desc) / 2 + (n - 1) * lh / 2 if va == "top" else y
        block = n * lh
        h = block + 8
        w = width if width is not None else max(
            ImageDraw.Draw(Image.new("RGB", (4, 4))).textlength(t, font=f) for t in lines) + 12
        tb = self._host(host).add_textbox(Pt(x), Pt(cy - h / 2), Pt(w), Pt(h))
        tf = tb.text_frame
        tf.word_wrap = wrap_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for side in ("left", "right", "top", "bottom"):
            setattr(tf, f"margin_{side}", 0)
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = Pt(lh)
            run = p.add_run()
            run.text = ln
            run.font.size = Pt(size * self.T.scale)
            run.font.name = self.font_name
            run.font.bold = weight in BOLD_WEIGHTS
            run.font.color.rgb = RGBColor(*color)
        if name:
            tb.name = name
        return tb


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("layout")
    ap.add_argument("--font", required=True, help="Inter.ttf, used for metrics")
    ap.add_argument("--out", required=True)
    ap.add_argument("--font-name", default="Inter", help="family PowerPoint asks for")
    ap.add_argument("--frames-root", default=None,
                    help="where the strip's frames live on this machine (see layout_common.py)")
    ap.add_argument("--type-scale", type=float, default=1.0)
    ap.add_argument("--no-group", action="store_true",
                    help="leave every token box a loose shape instead of grouping per frame")
    args = ap.parse_args()

    L = resolve_strip(json.load(open(args.layout)), args.layout, args.frames_root)
    T = Type(args.font, args.type_scale)
    d0 = ImageDraw.Draw(Image.new("RGB", (10, 10)))

    opts, gt = L["options"], L["answer"]
    letters = [chr(65 + i) for i in range(len(opts))]
    strip, rows = L["strip"], L["rows"]
    n = len(strip)

    # ── geometry, in design units, mirroring compose() ──
    dw = L["disp_w"]
    im0 = Image.open(strip[0]["path"])
    fw, fh = dw, int(dw * im0.height / im0.width)
    M, Lw, Rw, fgap, colgap = 46, 252, 340, 10, 26
    strip_w = n * fw + (n - 1) * fgap
    Wc = M + Lw + strip_w + colgap + Rw + M
    inner = Wc - 2 * M

    fq, fol, folb = T.pil(23, "SemiBold"), T.pil(16), T.pil(16, "SemiBold")
    fsub, fpred = T.pil(12), T.pil(15, "SemiBold")
    q_lh = int(fq.getmetrics()[0] * 1.34)
    opt_lh = int(fol.getmetrics()[0] * 1.42)
    sub_lh = int(fsub.getmetrics()[0] * 1.32)
    pred_lh = int(fpred.getmetrics()[0] * 1.42)

    qlines = wrap(d0, L["question"], fq, inner)
    blocks = [(wrap(d0, o.strip(), folb if letters[i] == gt else fol, inner - 16),
               letters[i] == gt) for i, o in enumerate(opts)]
    opts_h = sum(len(b) * opt_lh for b, _ in blocks) + (len(blocks) - 1) * 9
    y_q = M
    y_opts = y_q + len(qlines) * q_lh + 16
    y_legend = y_opts + opts_h + 18
    head = y_legend + 22 + 20

    sub_lines = [wrap(d0, r["sub"], fsub, Lw - 48) for r in rows]
    label_h = 28 + 6 + max(len(s) for s in sub_lines) * sub_lh + 10 + 18
    pred_texts = [(opts[letters.index(r["pred_letter"])].strip()
                   if r["pred_letter"] in letters else r["pred_letter"]) for r in rows]
    pred_lines = [wrap(d0, p, fpred, Rw - 8) for p in pred_texts]
    pred_h = 20 + max(len(p) for p in pred_lines) * pred_lh
    band_h = max(fh, pred_h, label_h + 20) + 32
    bands_end = head + len(rows) * band_h + (len(rows) - 1) * 12
    Hc = bands_end + 24 + M

    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(Wc), Pt(Hc)
    S = Slide(prs, T, args.font_name)

    # ── header ──
    S.text(M, M - 26, eyebrow(L), 12, "SemiBold", BLUE, name="eyebrow")
    S.text(M, y_q, qlines, 23, "SemiBold", INK, width=inner, line_h=q_lh, name="question")

    y = y_opts
    for i, (lines, is_gt) in enumerate(blocks):
        blk = len(lines) * opt_lh
        if is_gt:
            S.rect(M, y + 2, M + 4, y + blk - 4, fill=BLUE, line=None, name="gt-bar")
        S.text(M + 16, y, lines, 16, "SemiBold" if is_gt else "Regular",
               BLUE if is_gt else INK_80, width=inner - 16, line_h=opt_lh,
               name=f"option-{letters[i]}")
        y += blk + 9

    lx, cy = M, y_legend + 10.5
    for rgb, lab, is_ring in legend_entries(L):
        if is_ring:
            S.ring(lx + 7.5, cy, 7.5, rgb, 2.5, name="legend-gaze-ring")
        else:
            S.rrect(lx, y_legend + 3, lx + 15, y_legend + 18, 4, fill=rgb, line=None,
                    name="legend-swatch")
        S.text(lx + 22, cy, lab, 17, "Medium", INK, va="center", name="legend-label")
        lx += 22 + d0.textlength(lab, font=T.pil(17, "Medium")) + 26

    # ── rows ──
    s_h, s_w = L["grid"]
    by = head
    for ri, row in enumerate(rows):
        S.rrect(M - 10, by, Wc - M + 10, by + band_h, 16, fill=WHITE, line=HAIRLINE, lw=0.8,
                name=f"band-{row['name']}")
        vc = OK_FG if row["correct"] else BAD_FG
        cx0, cy0, cx1, cy1 = M, by + 12, M + Lw - 18, by + band_h - 12
        S.rrect(cx0, cy0, cx1, cy1, 12,
                fill=tuple(int(c * 0.16 + 255 * 0.84) for c in vc),
                line=tuple(int(c * 0.55 + 255 * 0.45) for c in vc), lw=0.8,
                name=f"card-{row['name']}")
        blk = 28 + 6 + len(sub_lines[ri]) * sub_lh + 10 + 18
        ty = cy0 + ((cy1 - cy0) - blk) / 2
        S.text(cx0 + 15, ty, row["name"], 20, "SemiBold", BLUE if row["is_ours"] else INK,
               width=Lw - 48, name=f"name-{row['name']}")
        S.text(cx0 + 15, ty + 34, sub_lines[ri], 12, "Regular", INK_80,
               width=Lw - 48, line_h=sub_lh, name=f"sub-{row['name']}")
        S.text(cx0 + 15, ty + 34 + len(sub_lines[ri]) * sub_lh + 6,
               "Correct" if row["correct"] else "Wrong", 12, "SemiBold", vc,
               width=Lw - 48, name=f"verdict-{row['name']}")

        fy = by + (band_h - fh) / 2
        fx = M + Lw
        # a row may override the shared strip with its own frame files (see
        # layout_common.resolve_strip); geometry is identical, only the pixels differ
        for si, cell in enumerate(row.get("strip") or strip):
            host = S.shapes if args.no_group else S.group(f"{row['name']}-t{cell['t']}")
            hs = host if args.no_group else host.shapes
            S.image(cell["path"], fx, fy, fw, fh, host=hs,
                    name=f"frame-{row['name']}-t{cell['t']}")
            pw, ph = fw / s_w, fh / s_h
            for grp in row["groups"][si]:
                rgb, lw = tuple(grp["rgb"]), max(0.6, grp["width_u"])
                a = grp["fill_alpha"] / 255 if grp["fill_alpha"] else 0
                for r, c in grp["cells"]:
                    S.rrect(fx + c * pw + 1, fy + r * ph + 1,
                            fx + (c + 1) * pw - 1, fy + (r + 1) * ph - 1,
                            max(1.0, pw * 0.18),
                            fill=(rgb if a else None), fill_alpha=a, line=rgb, lw=lw,
                            host=hs, name="token")
            if cell["gaze"]:
                gx, gy = fx + cell["gaze"][0] * fw, fy + cell["gaze"][1] * fh
                S.ring(gx, gy, 11, WHITE, 3, alpha=0.82, host=hs, name="gaze-halo")
                S.ring(gx, gy, 9, C_GAZE, 3, host=hs, name="gaze-ring")
            lbl = f"t{cell['t']}" + (" · final fixation" if cell["query_moment"] else "")
            tw = d0.textlength(lbl, font=T.pil(11, "SemiBold"))
            S.rrect(fx + 7, fy + fh - 29, fx + 7 + tw + 14, fy + fh - 7, 9,
                    fill=(BLUE if cell["query_moment"] else (0, 0, 0)),
                    fill_alpha=0.92 if cell["query_moment"] else 0.55,
                    line=None, host=hs, name="chip")
            S.text(fx + 14, fy + fh - 18, lbl, 11, "SemiBold", WHITE, va="center",
                   width=tw + 6, host=hs, name="chip-label", wrap_text=False)
            fx += fw + fgap

        px = fx + colgap
        pblk = 20 + len(pred_lines[ri]) * pred_lh
        py = by + (band_h - pblk) / 2
        S.text(px, py, "PREDICTION", 12, "SemiBold", INK_48, width=Rw - 8,
               name=f"predlabel-{row['name']}")
        S.text(px, py + 20, pred_lines[ri], 15, "SemiBold", OK_FG if row["correct"] else BAD_FG,
               width=Rw - 8, line_h=pred_lh, name=f"pred-{row['name']}")
        by += band_h + 12

    S.text(M, bands_end + 14, footer(L), 12, "Regular", INK_48, width=inner, name="footer")

    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    prs.save(args.out)
    nshapes = len(S.slide.shapes)
    print(f"-> {args.out}  ({Wc / 72:.1f} x {Hc / 72:.1f} in slide, {nshapes} top-level shapes, "
          f"{os.path.getsize(args.out) / 1e6:.1f} MB)")


if __name__ == "__main__":
    main()
